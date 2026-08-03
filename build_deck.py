# 15-slide pitch deck, Bol Impact-Report style. Open Sans, rounded shapes, green/blue palette,
# white + #0000a4 text, contrasting image slots. No em dashes; own insights; Bol slogan.
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

LGREEN = RGBColor(0x81,0xD2,0xA8); DGREEN = RGBColor(0x12,0x6B,0x47)
LBLUE  = RGBColor(0x9E,0xBD,0xF4); MBLUE  = RGBColor(0x67,0x97,0xE5); DBLUE = RGBColor(0x00,0x00,0xA4)
WHITE  = RGBColor(0xFF,0xFF,0xFF); YELLOW = RGBColor(0xFF,0xE5,0x00); RED = RGBColor(0xC0,0x39,0x2B)
FONT = "Open Sans"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    return s

def rrect(s, x, y, w, h, color, radius=0.09, line=None, lw=1.25):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try: sp.adjustments[0] = radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, size=18, color=DBLUE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, space=6, line=1.06):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [runs]
    for i, item in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.line_spacing = line
        text, ov = item if isinstance(item, tuple) else (item, {})
        r = p.add_run(); r.text = text
        f = r.font; f.name = FONT; f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold); f.color.rgb = ov.get("color", color)
    return tb

def stat(s, x, y, w, h, big, label, block, fg, sub=DBLUE):
    rrect(s, x, y, w, h, block, radius=0.12)
    txt(s, x, y+0.25, w, 1.0, big, size=48, color=fg, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+0.2, y+h-1.0, w-0.4, 0.9, label, size=14, color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.1)

def image_slot(s, x, y, w, h, block, hint):
    rrect(s, x, y, w, h, block, radius=0.07)
    rrect(s, x+0.35, y+0.35, 1.0, 1.0, WHITE, radius=0.18)
    txt(s, x+0.45, y+0.45, 0.8, 0.8, "img", size=18, color=block, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, y+h-0.85, w, 0.6, hint, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def head(s, kick, title, kcolor, tcolor, tw=7.2):
    rrect(s, 0.7, 0.6, 0.28, 0.46, kcolor, radius=0.4)
    txt(s, 1.12, 0.58, 11, 0.5, kick, size=13, color=kcolor, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if title: txt(s, 0.7, 1.15, tw, 1.7, title, size=32, color=tcolor, bold=True, line=1.0)

def footer(s, n, fg):
    txt(s, 0.7, 7.02, 9, 0.35, "Bol Sustainability Layer  ·  HACK4HER 2026", size=10, color=fg, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 12.0, 7.02, 0.8, 0.35, str(n), size=10, color=fg, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# 1. TITLE
s = slide(DBLUE)
rrect(s, 7.7, 1.1, 5.0, 5.3, MBLUE, radius=0.06)
txt(s, 8.1, 1.5, 4.2, 0.5, "REAL PRODUCT, REAL DATA", size=12, color=YELLOW, bold=True)
txt(s, 8.1, 2.05, 4.3, 3.6,
    [("Nano Phone Speaker", {"size":22,"color":WHITE,"bold":True}),
     ("Cradle to Cradle certified.", {"size":17,"color":LBLUE}),
     ("Shipped in plastic, from Japan.", {"size":17,"color":LBLUE}),
     ("168 kg of CO2 to make.", {"size":17,"color":LBLUE}),
     ("Our verdict: 2 leaves, badge withheld.", {"size":17,"color":YELLOW,"bold":True})], line=1.25, space=10)
rrect(s, 0.7, 2.0, 0.6, 0.16, YELLOW, radius=0.4)
txt(s, 0.7, 2.25, 6.8, 0.6, "HACK4HER × BOL · TRACK 2: THE BRAIN", size=14, color=YELLOW, bold=True)
txt(s, 0.65, 2.85, 7.0, 1.9, "The Honest Price Tag", size=52, color=WHITE, bold=True, line=0.95)
txt(s, 0.7, 4.7, 6.8, 1.5,
    "We score how sustainable a product really is, show it in one glance, and quietly flag the sellers whose green claims do not hold up.",
    size=19, color=LBLUE, line=1.2)
txt(s, 0.7, 6.5, 7, 0.5, "Team [your names]  ·  bol-com-solution.vercel.app", size=12, color=LBLUE)

# 2. PROBLEM
s = slide(LGREEN)
head(s, "THE PROBLEM", "People want to shop greener. The shop makes it hard.", DGREEN, DBLUE)
image_slot(s, 8.3, 1.2, 4.4, 5.2, DGREEN, "PHOTO: shopper at checkout")
txt(s, 0.7, 2.7, 7.3, 3.4,
    [("Sustainability information is scattered across the page, written in jargon, or simply not there.", {"size":21,"color":DBLUE,"bold":True}),
     ("So a shopper who cares ends up doing the easy thing: the cheapest item with the fastest delivery. "
      "The good intention is there. The information to act on it is not.", {"size":18,"color":DGREEN})], line=1.2, space=14)
footer(s, 2, DGREEN)

# 3. BOL SLOGAN
s = slide(DBLUE)
rrect(s, 0.7, 1.6, 0.6, 0.16, YELLOW, radius=0.4)
txt(s, 0.7, 1.85, 11, 0.5, "BOL'S OWN PROMISE", size=14, color=YELLOW, bold=True)
txt(s, 0.65, 2.45, 12.0, 1.4, "“de winkel van ons allemaal”", size=48, color=WHITE, bold=True, line=1.0)
txt(s, 0.7, 3.75, 12, 0.6, "the shop that belongs to all of us", size=20, color=LBLUE)
txt(s, 0.7, 4.7, 12.0, 1.6,
    [("Bol says it wants gemak, betaalbaarheid en duurzaamheid, convenience, affordability and sustainability, to go together.", {"size":19,"color":WHITE,"bold":True}),
     ("Right now shoppers feel they have to pick. Our layer is built to put all three on the same product.", {"size":18,"color":LBLUE})], line=1.2, space=12)
footer(s, 3, LBLUE)

# 4. SOLUTION
s = slide(LBLUE)
head(s, "OUR ANSWER", "A thin sustainability layer inside the shop", DBLUE, DBLUE)
image_slot(s, 8.3, 1.2, 4.4, 5.2, DBLUE, "SCREENSHOT: product card + leaf score")
cards = [("One score, as leaves", "Right next to the price."),
         ("A reason you can read", "Plain language, never a black box."),
         ("A greener alternative", "In the same category, usually cheaper."),
         ("A check on the claim", "We notice when a label does not add up.")]
y = 2.5
for hgt, sub in cards:
    rrect(s, 0.7, y, 7.3, 0.92, WHITE, radius=0.2)
    txt(s, 1.0, y+0.1, 3.2, 0.72, hgt, size=16, color=DBLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 4.2, y+0.1, 3.6, 0.72, sub, size=14, color=MBLUE, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02
footer(s, 4, DBLUE)

# 5. HOW THE SCORE WORKS
s = slide(MBLUE)
head(s, "THE BRAIN", "One score, four stages of a product's life", WHITE, WHITE, tw=11)
cols = [("MAKE","Manufacturing CO2","35%"),("MOVE","Origin to NL","15%"),
        ("PACK","Packaging and warehouse","25%"),("LAST","Repairability and lifespan","25%")]
x = 0.7
for name, desc, w in cols:
    rrect(s, x, 2.4, 2.95, 2.3, DBLUE, radius=0.1)
    txt(s, x+0.25, 2.65, 2.5, 0.6, name, size=22, color=YELLOW, bold=True)
    txt(s, x+0.25, 3.3, 2.5, 1.0, desc, size=14, color=WHITE, line=1.12)
    txt(s, x+0.25, 4.25, 2.5, 0.4, w+" of the score", size=12, color=LBLUE, bold=True)
    x += 3.07
txt(s, 0.7, 5.05, 12.0, 1.6,
    [("Each stage is ranked against products in the same category, then blended into one number from 0 to 100.", {"size":16,"color":WHITE,"bold":True}),
     ("The weights follow the EU Product Environmental Footprint method, where climate change carries the most weight of any single category.", {"size":15,"color":LBLUE})],
    line=1.18, space=8)
footer(s, 5, LBLUE)

# 6. INSIGHT: CO2 IS MANUFACTURING
s = slide(LGREEN)
head(s, "WHAT WE FOUND IN THE DATA", "The footprint is about making, not shipping", DGREEN, DBLUE, tw=8)
stat(s, 9.0, 2.4, 3.7, 3.4, "0.12", "correlation between shipping distance and CO2", DGREEN, WHITE)
txt(s, 0.7, 2.5, 7.9, 3.6,
    [("Before we trusted the carbon field, we tested it.", {"size":20,"color":DBLUE,"bold":True}),
     ("Across the whole catalog, CO2 barely moves with how far a product ships (correlation 0.12), and every digital product reads zero. "
      "That is the fingerprint of a manufacturing footprint.", {"size":18,"color":DGREEN}),
     ("So Make and Move measure genuinely different things, and no gram of carbon is counted twice.", {"size":18,"color":DBLUE,"bold":True})],
    line=1.2, space=12)
footer(s, 6, DGREEN)

# 7. INSIGHT: MISSING DATA
s = slide(LBLUE)
head(s, "HANDLING THE GAPS HONESTLY", "Half the sustainability data is missing. We say so.", DBLUE, DBLUE, tw=8)
y = 2.45
for val, lab in [("52%","recyclability blank"),("49%","lifespan blank"),("26%","carbon footprint blank")]:
    rrect(s, 9.0, y, 3.7, 0.95, WHITE, radius=0.2)
    txt(s, 9.2, y+0.1, 1.5, 0.75, val, size=28, color=DBLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 10.7, y+0.1, 1.9, 0.75, lab, size=14, color=MBLUE, anchor=MSO_ANCHOR.MIDDLE, line=1.05)
    y += 1.08
txt(s, 0.7, 2.5, 7.9, 3.7,
    [("Rather than hide the gaps, every score carries a confidence level set by how complete its data is.", {"size":19,"color":DBLUE,"bold":True}),
     ("Low confidence scores fade and say so on the card.", {"size":18,"color":MBLUE}),
     ("The same number doubles as a readiness check for the EU Digital Product Passport, which will soon require these exact fields.", {"size":18,"color":DBLUE})],
    line=1.2, space=12)
footer(s, 7, DBLUE)

# 8. GREENWASHING
s = slide(LGREEN)
head(s, "THE CHECK", "We compare the claim to the evidence", DGREEN, DBLUE, tw=8)
image_slot(s, 8.6, 1.3, 4.1, 5.0, DGREEN, "SCREENSHOT: a flagged product")
rrect(s, 0.7, 2.4, 7.5, 1.5, WHITE, radius=0.14, line=RED, lw=1.5)
txt(s, 1.0, 2.6, 7.0, 1.1,
    "Carbon neutral badge, on an item in the dirtiest quarter of its category for CO2.",
    size=18, color=RED, bold=True, anchor=MSO_ANCHOR.MIDDLE, line=1.15)
txt(s, 0.7, 4.2, 7.6, 2.2,
    [("A small rule engine looks for contradictions like that one, plus eco labels on plastic, far shipped items, "
      "and recyclable claims on barely recyclable material.", {"size":17,"color":DGREEN}),
     ("When the claim and the data disagree, the badge quietly stays off. We inform the shopper, we do not shame the seller.", {"size":17,"color":DBLUE,"bold":True})],
    line=1.2, space=10)
footer(s, 8, DGREEN)

# 9. INSIGHT: GREENER ISN'T PRICIER
s = slide(LBLUE)
head(s, "THE MYTH WE CAN BUST", "Greener almost never costs more", DBLUE, DBLUE, tw=8)
stat(s, 9.0, 2.3, 3.7, 1.6, "0.03", "price vs sustainability correlation", MBLUE, WHITE)
stat(s, 9.0, 4.1, 3.7, 1.6, "81%", "products with a cheaper greener twin", DBLUE, WHITE)
txt(s, 0.7, 2.45, 7.9, 3.8,
    [("In Bol's catalog, price and sustainability are essentially unrelated.", {"size":20,"color":DBLUE,"bold":True}),
     ("For 81% of products there is a greener option in the same category that also costs less. "
      "That is why our nudge leads with the saving, not with guilt.", {"size":18,"color":MBLUE}),
     ("Convenience, affordability and sustainability really can sit on the same product.", {"size":18,"color":DBLUE,"bold":True})],
    line=1.2, space=12)
footer(s, 9, DBLUE)

# 10. CHECKOUT AGGREGATE
s = slide(LGREEN)
head(s, "MAKE THE IMPACT REAL", "Small swaps barely register. A basket does.", DGREEN, DBLUE, tw=8)
image_slot(s, 8.6, 1.3, 4.1, 5.0, DGREEN, "SCREENSHOT: basket savings line")
txt(s, 0.7, 2.5, 7.6, 3.7,
    [("Per product the saving is tiny. Across a basket it adds up.", {"size":20,"color":DBLUE,"bold":True}),
     ("At checkout we total the money and the CO2 saved, and translate the carbon into something you can picture: "
      "a drive to Rotterdam, a drive to Berlin, a short flight to London.", {"size":18,"color":DGREEN}),
     ("Car at 170 g per km; flights from standard per passenger figures.", {"size":14,"color":DGREEN})],
    line=1.2, space=12)
footer(s, 10, DGREEN)

# 11. DEMO
s = slide(DBLUE)
image_slot(s, 7.6, 1.1, 5.1, 5.3, MBLUE, "LIVE: bol-com-solution.vercel.app")
rrect(s, 0.7, 2.0, 0.6, 0.16, YELLOW, radius=0.4)
txt(s, 0.7, 2.25, 6.6, 0.6, "LIVE DEMO", size=15, color=YELLOW, bold=True)
txt(s, 0.65, 2.85, 6.8, 1.3, "Let's shop.", size=54, color=WHITE, bold=True)
txt(s, 0.7, 4.3, 6.7, 2.0,
    "Browse the shop, glance the score, open the breakdown, see a greener and cheaper pick, then watch the basket total your impact at checkout.",
    size=19, color=LBLUE, line=1.2)
footer(s, 11, LBLUE)

# 12. BUILT ON BOL
s = slide(LBLUE)
head(s, "COULD BOL REALLY SHIP IT", "It fits what Bol already runs", DBLUE, DBLUE, tw=8)
image_slot(s, 8.6, 1.3, 4.1, 5.0, DGREEN, "PHOTO: Goede Keuze / fulfilment")
pts = ["Bol already labels a million products Goede Keuze. That label is a yes or no. We make it a ranked, explainable score.",
       "We reuse Bol's own rule for it: a recognised certificate, or at least 50% recycled material.",
       "We lean on Bol's real figures. Fit to size packing saves 28% of packaging CO2 per parcel.",
       "Transport is only 1.5% of Bol's emissions, which is exactly why we put most of the weight on manufacturing."]
y = 2.45
for p in pts:
    rrect(s, 0.75, y+0.07, 0.18, 0.18, MBLUE, radius=0.4)
    txt(s, 1.15, y-0.05, 7.0, 0.95, p, size=15, color=DBLUE, line=1.13)
    y += 0.93
footer(s, 12, DBLUE)

# 13. WHAT'S NEXT
s = slide(LGREEN)
head(s, "IF WE HAD ANOTHER MONTH", "Score the second life, and the scale", DGREEN, DBLUE, tw=8)
image_slot(s, 8.6, 1.3, 4.1, 5.0, DGREEN, "ICON: refurb, repair, recycle")
txt(s, 0.7, 2.5, 7.6, 3.8,
    [("Bol already refurbishes, repairs and resells. We would feed that back into the score: "
      "is a refurbished version on offer, can this be repaired through Bol.", {"size":18,"color":DGREEN}),
     ("We would also rank products by total impact, units sold times footprint, so Bol can see where a nudge saves tonnes, not grams. "
      "Our catalog already adds up to roughly 2.9 million tonnes of embodied CO2.", {"size":18,"color":DBLUE,"bold":True})],
    line=1.2, space=14)
footer(s, 13, DGREEN)

# 14. CLOSE / SLOGAN
s = slide(DGREEN)
rrect(s, 0.7, 1.9, 0.6, 0.16, LGREEN, radius=0.4)
txt(s, 0.7, 2.15, 11, 0.5, "WHERE WE WANT TO TAKE IT", size=14, color=LGREEN, bold=True)
txt(s, 0.65, 2.8, 12.2, 3.0,
    [("A winkel van ons allemaal is also a shop that is honest with all of us:", {"size":34,"color":WHITE,"bold":True}),
     ("one glance, the real story, and the greener choice that happens to cost less.", {"size":34,"color":LGREEN,"bold":True})],
    line=1.1, space=10)
footer(s, 14, LGREEN)

# 15. THANK YOU
s = slide(DBLUE)
rrect(s, 8.3, 1.4, 4.4, 4.7, MBLUE, radius=0.07)
txt(s, 8.6, 1.75, 3.85, 4.1,
    [("CO2 is manufacturing only. We tested it (correlation 0.12).", {"size":15,"color":WHITE}),
     ("Greener and cheaper on 4 of 5 products. No premium asked.", {"size":15,"color":WHITE}),
     ("Built for the EU Green Claims rules that are already coming.", {"size":15,"color":WHITE})],
    line=1.2, space=14)
txt(s, 0.7, 2.6, 7.2, 1.3, "Thank you.", size=56, color=WHITE, bold=True)
txt(s, 0.7, 4.0, 7.2, 0.7, "Questions?", size=26, color=YELLOW, bold=True)
txt(s, 0.7, 5.0, 7.0, 0.5, "Team [your names]  ·  bol-com-solution.vercel.app", size=13, color=LBLUE)

prs.save("Bol_Sustainability_Layer_Pitch.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")
